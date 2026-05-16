"""Generate DecayMonitor pitch deck — Apple Keynote style (dark theme)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

# ── Color Palette (Apple-style dark) ──────────────────────────────────
BG_DARK   = RGBColor(0x0A, 0x0A, 0x0F)   # near-black
BG_CARD   = RGBColor(0x1C, 0x1C, 0x24)   # dark card
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_L    = RGBColor(0xA1, 0xA1, 0xAA)   # secondary text
GRAY_M    = RGBColor(0x6B, 0x6B, 0x76)   # dim text
ACCENT    = RGBColor(0x5E, 0x5C, 0xE6)   # indigo accent
ACCENT2   = RGBColor(0x3B, 0x82, 0xF6)   # blue
ACCENT3   = RGBColor(0xF5, 0x9E, 0x0B)   # amber (warnings)
GREEN     = RGBColor(0x34, 0xD3, 0x99)   # success
RED       = RGBColor(0xEF, 0x44, 0x44)   # danger
PURPLE    = RGBColor(0xA7, 0x8B, 0xFA)   # soft purple

SLIDE_W = Inches(13.333)  # 16:9
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

# ── Helpers ───────────────────────────────────────────────────────────

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Helvetica Neue", line_spacing=1.1):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    # line spacing
    pPr = p._pPr
    if pPr is None:
        pPr = p._p.get_or_add_pPr()
    lnSpc = pPr.makeelement(qn('a:lnSpc'), {})
    spcPct = lnSpc.makeelement(qn('a:spcPct'), {'val': str(int(line_spacing * 100000))})
    lnSpc.append(spcPct)
    pPr.append(lnSpc)
    return tf

def add_multiline(slide, left, top, width, height, lines, font_size=16,
                  color=WHITE, bold_first=False, font_name="Helvetica Neue",
                  line_spacing=1.3, alignment=PP_ALIGN.LEFT):
    """Add textbox with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.font.bold = (bold_first and i == 0)
        p.alignment = alignment
        p.space_after = Pt(4)
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        lnSpc = pPr.makeelement(qn('a:lnSpc'), {})
        spcPct = lnSpc.makeelement(qn('a:spcPct'), {'val': str(int(line_spacing * 100000))})
        lnSpc.append(spcPct)
        pPr.append(lnSpc)
    return tf

def add_rounded_rect(slide, left, top, width, height, fill_color=BG_CARD,
                     border_color=None, corner_radius=8):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    # Adjust corner radius via XML
    sp = shape._element
    prstGeom = sp.find(qn('a:prstGeom'), sp.nsmap) if hasattr(sp, 'nsmap') else None
    if prstGeom is None:
        prstGeom = sp.find('.//' + qn('a:prstGeom'))
    return shape

def add_circle(slide, left, top, diameter, fill_color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left), Inches(top), Inches(diameter), Inches(diameter)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_gradient_bg(slide, colors):
    """Add gradient background via XML. colors = list of (RGBColor, position)."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = colors[0][0]  # fallback to first color

def add_big_number(slide, left, top, number, label, num_color=WHITE, label_color=GRAY_L):
    tf = add_textbox(slide, left, top, 4, 1.2, str(number), font_size=64,
                     color=num_color, bold=True)
    add_textbox(slide, left, top + 1.1, 4, 0.6, label, font_size=14, color=label_color)

def add_page_number(slide, num):
    add_textbox(slide, 12.2, 7.0, 0.8, 0.4, str(num), font_size=10, color=GRAY_M,
                alignment=PP_ALIGN.RIGHT)

def add_section_header(slide, title, subtitle=""):
    add_textbox(slide, 0.8, 0.6, 8, 0.8, title, font_size=36, color=WHITE, bold=True)
    if subtitle:
        add_textbox(slide, 0.8, 1.3, 10, 0.5, subtitle, font_size=16, color=GRAY_L)

def add_thin_line(slide, left, top, width, color=GRAY_M):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Pt(1)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

# ── Slide 1: Title ────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, BG_DARK)

# Large central title
add_textbox(slide, 1.5, 1.8, 10.3, 1.5, "DecayMonitor", font_size=80,
            color=WHITE, bold=True)
add_textbox(slide, 1.5, 3.2, 10.3, 0.8, "LLM Recursive Stability  β", font_size=42,
            color=ACCENT, bold=True)
add_thin_line(slide, 1.5, 4.2, 3.5, ACCENT)
add_textbox(slide, 1.5, 4.6, 8, 0.6,
            "The first benchmark for LLM self-consuming stability",
            font_size=20, color=GRAY_L)
add_textbox(slide, 1.5, 5.5, 8, 0.5,
            "DecayMonitor.ai  ·  Beijing  ·  2026",
            font_size=14, color=GRAY_M)
add_page_number(slide, 1)

# ── Slide 2: The Problem ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_header(slide, "The Self-Consumption Loop", "Why recursive stability matters now")

# Left: the cycle diagram
add_textbox(slide, 0.8, 2.2, 6, 0.6, "Training Data → Model → Generated Content → Training Data → …",
            font_size=22, color=WHITE, bold=True)
add_textbox(slide, 0.8, 3.0, 6, 1.2,
            "Every generation adds noise.\n"
            "As LLM-generated content floods the internet,\n"
            "models increasingly train on their own outputs.",
            font_size=16, color=GRAY_L)

# Right: stats cards
add_rounded_rect(slide, 7.5, 2.2, 5.0, 1.3, BG_CARD)
add_textbox(slide, 7.8, 2.3, 4.5, 0.5, "57%", font_size=48, color=ACCENT, bold=True)
add_textbox(slide, 7.8, 2.9, 4.5, 0.4, "of web content will be AI-generated by 2027", font_size=13, color=GRAY_L)

add_rounded_rect(slide, 7.5, 3.7, 5.0, 1.3, BG_CARD)
add_textbox(slide, 7.8, 3.8, 4.5, 0.5, "Model Collapse", font_size=36, color=RED, bold=True)
add_textbox(slide, 7.8, 4.4, 4.5, 0.4, "Documented in Nature (Shumailov et al., 2024)", font_size=13, color=GRAY_L)

add_rounded_rect(slide, 7.5, 5.2, 5.0, 1.3, BG_CARD)
add_textbox(slide, 7.8, 5.3, 4.5, 0.5, "No Standard Metric", font_size=36, color=ACCENT2, bold=True)
add_textbox(slide, 7.8, 5.9, 4.5, 0.4, "The industry has no way to measure recursive stability", font_size=13, color=GRAY_L)

add_page_number(slide, 2)

# ── Slide 3: The Gap ──────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_header(slide, "The Evaluation Blindspot", "All existing benchmarks measure single-pass quality. No one measures what happens across generations.")

# Table-like layout
benchmarks = [
    ("MMLU", "Knowledge breadth", "Single-pass accuracy"),
    ("HumanEval", "Code generation", "Pass@k on static tests"),
    ("Chatbot Arena", "Human preference", "Elo from pairwise votes"),
    ("HELM", "Multi-dimensional", "Holistic single-generation"),
    ("DecayMonitor", "Recursive stability β", "Cross-generation constraint decay"),
]
y = 2.3
for name, desc, metric in benchmarks:
    is_us = (name == "DecayMonitor")
    c = ACCENT if is_us else BG_CARD
    add_rounded_rect(slide, 0.8, y, 11.5, 0.85, c, border_color=ACCENT if is_us else None)
    add_textbox(slide, 1.1, y + 0.1, 3.0, 0.5, name, font_size=18, color=WHITE if is_us else WHITE, bold=is_us)
    add_textbox(slide, 4.2, y + 0.1, 3.5, 0.5, desc, font_size=14, color=GRAY_L)
    add_textbox(slide, 7.8, y + 0.1, 4.2, 0.5, metric, font_size=14, color=ACCENT2 if is_us else GRAY_M)
    y += 1.0

add_page_number(slide, 3)

# ── Slide 4: Our Solution — β ─────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_header(slide, "β: The Recursive Stability Index")

# Big beta
add_circle(slide, 4.8, 1.8, 2.6, ACCENT)
add_textbox(slide, 5.5, 2.5, 1.6, 1.4, "β", font_size=72, color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, 0.8, 4.8, 11.5, 0.8,
            "β measures how fast LLM-generated content degrades when models iterate on their own outputs.",
            font_size=18, color=WHITE, bold=True)

add_multiline(slide, 0.8, 5.6, 11.5, 1.5, [
    "Sₙ = Sₙ₋₁ · (1 − β)    —    Recursive stability decay",
    "β ∈ [0.001, 0.55]    —    Lower β = more stable across generations",
    "Measured via constraint residual total_constraint = Σ|∇σᵢ|   —    No LLM judge dependency",
], font_size=16, color=GRAY_L, line_spacing=1.5)

add_page_number(slide, 4)

# ── Slide 5: Algorithm Pipeline ───────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_header(slide, "Algorithm Pipeline", "CPU-only, rule-based, no model dependency — $0.50 per evaluation")

pipeline = [
    ("1", "Seed Text\n(Gen 0)", GREEN),
    ("2", "API Generate\nGen 1→2→3", ACCENT2),
    ("3", "8 Text Features\nRule-based extraction", ACCENT),
    ("4", "5D σ Mapping\nσ_fact, σ_syntax, σ_style", PURPLE),
    ("5", "Π = Σ∇σᵢ\nConstraint residual", RGBColor(0xF5, 0x9E, 0x0B)),
    ("6", "Exponential Fit\nβ = 1 − eˢˡᵒᵖᵉ", RED),
]

x = 0.4
for num, desc, color in pipeline:
    add_rounded_rect(slide, x, 2.2, 1.85, 3.5, BG_CARD)
    add_circle(slide, x + 0.6, 2.5, 0.65, color)
    add_textbox(slide, x + 0.6, 2.6, 0.65, 0.5, num, font_size=22, color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.1, 3.5, 1.65, 1.8, desc, font_size=13, color=GRAY_L,
                alignment=PP_ALIGN.CENTER)
    # arrow between boxes (except last)
    if num != "6":
        add_textbox(slide, x + 1.85, 3.4, 0.35, 0.5, "→", font_size=24, color=GRAY_M,
                    alignment=PP_ALIGN.CENTER)
    x += 2.08

add_textbox(slide, 0.8, 6.2, 11.5, 0.6,
            "Fully reproducible  ·  CV = 3.3%  ·  Pre-registered model (no post-hoc selection bias)",
            font_size=14, color=GRAY_M)

add_page_number(slide, 5)

# ── Slide 6: 9-Model β Ranking ────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_header(slide, "9-Model β Ranking", "Pre-registered exponential + total_constraint  ·  n = 100 seeds per model")

models = [
    ("DeepSeek-Chat (V3)", 0.0281, 1, GREEN),
    ("GPT-4o-mini", 0.0885, 2, WHITE),
    ("Llama 3.1 70B", 0.0925, 3, WHITE),
    ("Llama 3.1 8B", 0.0942, 4, WHITE),
    ("GPT-4o", 0.0985, 5, WHITE),
    ("DeepSeek-R1", 0.1038, 6, WHITE),
    ("Claude Sonnet 4.6", 0.1055, 7, WHITE),
    ("Claude Opus 4.6", 0.1196, 8, WHITE),
    ("Claude Haiku 4.5", 0.1468, 9, RED),
]

# Bar chart
y_start = 2.2
bar_area_w = 8.0
max_width = 7.5  # inches for max beta
max_beta = 0.16

for i, (name, beta, rank, color) in enumerate(models):
    y = y_start + i * 0.55
    # label
    add_textbox(slide, 0.8, y, 3.3, 0.4, name, font_size=14, color=WHITE)
    # bar background
    add_rounded_rect(slide, 4.3, y + 0.02, bar_area_w, 0.35, RGBColor(0x2A, 0x2A, 0x35))
    # bar fill
    bar_w = (beta / max_beta) * bar_area_w
    c = color
    if beta == 0.0281:
        c = GREEN
    elif beta <= 0.11:
        c = ACCENT2
    else:
        c = RED
    add_rounded_rect(slide, 4.3, y + 0.02, bar_w, 0.35, c)
    # value
    add_textbox(slide, 4.3 + bar_w + 0.15, y + 0.02, 1.2, 0.35, f"{beta:.4f}",
                font_size=14, color=WHITE, bold=True)

# Insights
insights = [
    "DeepSeek-V3: 3.2x better than #2 — clearly separated leader",
    "Mid-tier (2−7): β ∈ [0.089, 0.106] — statistically indistinguishable",
    "Model family β is constant within families (OpenAI range=0.010, Llama range=0.002)",
]
add_multiline(slide, 0.8, 5.8, 11.5, 1.5, insights, font_size=13, color=GRAY_L, line_spacing=1.4)
add_page_number(slide, 6)

# ── Slide 7: Capability Heatmap ───────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_header(slide, "Per-Capability Decay Heatmap (Gen3 Sₙ)", "creative_writing — the universal Achilles' heel")

caps = ["math", "code", "factual", "logic", "creative", "general"]
models_hm = [
    ("DeepSeek-Chat",   ["H","H","H","H","C","H"]),
    ("GPT-4o-mini",     ["D","H","H","H","C","C"]),
    ("GPT-4o",          ["D","H","H","H","X","D"]),
    ("Llama 3.1 70B",   ["D","H","H","C","X","D"]),
    ("Llama 3.1 8B",    ["X","C","H","H","X","X"]),
    ("DeepSeek-R1",     ["H","H","H","D","C","D"]),
    ("Claude Sonnet",   ["C","H","H","C","C","C"]),
    ("Claude Opus",     ["D","D","H","D","C","H"]),
    ("Claude Haiku",    ["C","D","H","D","C","C"]),
]

STATUS_COLORS = {
    "H": (GREEN, "Healthy >0.8"),
    "D": (ACCENT2, "Degrading 0.5−0.8"),
    "C": (ACCENT3, "Critical 0.3−0.5"),
    "X": (RED, "Collapsed <0.3"),
}

# Header row
x_start = 0.8
col_w = 1.55
row_h = 0.42
add_textbox(slide, x_start, 2.0, 2.0, 0.4, "", font_size=12)  # empty cell
for j, cap in enumerate(caps):
    add_textbox(slide, x_start + 2.0 + j * col_w, 2.0, col_w, 0.4, cap,
                font_size=10, color=GRAY_L, alignment=PP_ALIGN.CENTER)

for i, (model, statuses) in enumerate(models_hm):
    y = 2.4 + i * row_h
    add_textbox(slide, x_start, y, 2.0, 0.35, model, font_size=12, color=WHITE)
    for j, st in enumerate(statuses):
        c, _ = STATUS_COLORS[st]
        add_rounded_rect(slide, x_start + 2.0 + j * col_w + 0.08, y, col_w - 0.15, row_h - 0.08, c)
        add_textbox(slide, x_start + 2.0 + j * col_w + 0.08, y + 0.02, col_w - 0.15, row_h - 0.08,
                    st, font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Legend
legend_y = 2.4 + len(models_hm) * row_h + 0.3
legend_items = list(STATUS_COLORS.items())
for k, (st, (c, desc)) in enumerate(legend_items):
    lx = x_start + k * 2.8
    add_rounded_rect(slide, lx, legend_y, 0.35, 0.25, c)
    add_textbox(slide, lx + 0.45, legend_y - 0.02, 2.2, 0.3, desc, font_size=10, color=GRAY_L)

add_textbox(slide, 0.8, legend_y + 0.5, 11.5, 0.4,
            "9/9 models collapse or near-collapse on creative_writing  ·  factual_knowledge: 9/9 healthy",
            font_size=13, color=GRAY_L)
add_page_number(slide, 7)

# ── Slide 8: Scientific Rigor ─────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_header(slide, "Scientific Rigor", "Designed to withstand peer review")

# Four rigor cards
cards = [
    ("Test-Retest\nCV = 3.3%", "Perfect reproduction\nβ=0.0885 = 0.0885\nPost-hoc CV was 11.3%\n→ 3.4x improvement", ACCENT),
    ("Bootstrap\nConvergence", "β stabilizes at n≥60\nn=100 recommended\nCI width: 0.156\n30% narrower than n=36", ACCENT2),
    ("Seed\nSensitivity", "Δ = 0.005\nNegligible impact\nBelow measurement noise\nDifferent seed sets agree", GREEN),
    ("Pre-Registered\nMethod", "No model/target competition\nEliminates +184% inflation\nExponential + total_constraint\nHonest β for every model", PURPLE),
]

for i, (title, desc, color) in enumerate(cards):
    x = 0.8 + i * 3.1
    add_rounded_rect(slide, x, 2.2, 2.85, 3.8, BG_CARD)
    add_textbox(slide, x + 0.2, 2.4, 2.5, 1.0, title, font_size=22, color=color, bold=True)
    add_thin_line(slide, x + 0.2, 3.3, 1.2, color)
    add_textbox(slide, x + 0.2, 3.5, 2.5, 2.3, desc, font_size=14, color=GRAY_L)

add_textbox(slide, 0.8, 6.4, 11.5, 0.5,
            "Also validated: P3 neural collapse (λ_C=+0.0415, R²=0.884)  ·  4-family coverage  ·  Temperature robustness",
            font_size=13, color=GRAY_M)
add_page_number(slide, 8)

# ── Slide 9: P3 Neural Validation ─────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_header(slide, "Neural Validation: Constraint Attractor Collapse",
                   "Independent verification on Qwen2.5-1.5B — constraint diversity decays exponentially")

add_rounded_rect(slide, 0.8, 2.2, 5.8, 4.5, BG_CARD)
add_textbox(slide, 1.1, 2.4, 5.3, 0.5, "P3 Rigorous Test", font_size=22, color=ACCENT, bold=True)
add_multiline(slide, 1.1, 3.0, 5.3, 3.5, [
    "C_div(n) = C_div(0) · e^(−λ_C · n)",
    "",
    "λ_C (decay rate)     = +0.0415  ✓",
    "R² (exponential fit) = 0.884    ✓",
    "Mean ||Π|| CV        = 0.301    ✓",
    "Gen3/Gen0 ratio      = 0.926    ✓",
    "",
    "Constraint diversity decays monotonically",
    "across generations — Attractor Collapse confirmed.",
], font_size=14, color=GRAY_L, line_spacing=1.3)

# Right: key finding
add_rounded_rect(slide, 7.2, 2.2, 5.3, 4.5, BG_CARD)
add_textbox(slide, 7.5, 2.4, 4.8, 0.5, "What This Means", font_size=22, color=GREEN, bold=True)
add_multiline(slide, 7.5, 3.0, 4.8, 3.5, [
    "β is not just a statistical artifact —",
    "it has a neural correlate.",
    "",
    "As models self-consume:",
    "• Constraint attractors collapse",
    "• Generation patterns converge",
    "• Response diversity decreases",
    "",
    "This is the mechanism behind β.",
    "It's grounded in neural dynamics.",
], font_size=14, color=GRAY_L, line_spacing=1.3)

add_page_number(slide, 9)

# ── Slide 10: Business Model ──────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_header(slide, "Business Model", "Free leaderboard + paid enterprise services")

tiers = [
    ("Free Leaderboard", "$0", "Public β rankings\n6-capability heatmaps\nModel detail pages\nCross-model comparison\nWeekly auto-refresh", GREEN),
    ("Pro — Private Eval", "$1,499/eval", "Pre-release testing\nResults stay off leaderboard\nE-I/E-II/E-III deep diagnosis\nAnonymous competitor comparison\nPDF + CSV/JSONL raw data", ACCENT2),
    ("Enterprise", "$4,999/mo", "Per-checkpoint β tracking\nTraining data blindspot diagnosis\nDedicated API endpoint\nSlack integration\nQuarterly strategy calls", PURPLE),
]

for i, (name, price, features, color) in enumerate(tiers):
    x = 0.8 + i * 4.1
    add_rounded_rect(slide, x, 2.2, 3.85, 4.5, BG_CARD, border_color=color if i > 0 else None)
    add_textbox(slide, x + 0.3, 2.4, 3.3, 0.5, name, font_size=20, color=color, bold=True)
    add_textbox(slide, x + 0.3, 2.9, 3.3, 0.5, price, font_size=32, color=WHITE, bold=True)
    add_thin_line(slide, x + 0.3, 3.5, 1.5, color)
    add_textbox(slide, x + 0.3, 3.7, 3.3, 2.7, features, font_size=14, color=GRAY_L)

# Market reference
add_textbox(slide, 0.8, 7.0, 11.5, 0.4,
            "Market reference: LMArena = $17B valuation, $30M ARR  ·  We add a new dimension to the LLM evaluation matrix",
            font_size=12, color=GRAY_M)
add_page_number(slide, 10)

# ── Slide 11: Competitive Moat ────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_header(slide, "Competitive Moat", "Why DecayMonitor is defensible")

moats = [
    ("Definition\nOwnership", "If 'recursive stability' becomes a standard evaluation dimension,\nthe definer owns the category — like LMArena owns Elo for chatbots.", ACCENT),
    ("Full-Auto\nLow Cost", "$0.50/model evaluation. CPU-only rule-based extraction.\nNo GPU, no LLM judge, no 5M human votes needed.", ACCENT2),
    ("Cross-Family\nβ Database", "4+ model families benchmarked (OpenAI, Claude, Llama, DeepSeek).\nFirst-mover data advantage. Every new model adds to the moat.", GREEN),
    ("Theory\nMoat", "Constraint Attractor Collapse + P1−P4 testable predictions.\nNot just a metric — a causal theory with neural validation.", PURPLE),
]

for i, (title, desc, color) in enumerate(moats):
    y = 2.2 + i * 1.25
    add_rounded_rect(slide, 0.8, y, 11.5, 1.1, BG_CARD, border_color=color if i == 0 else None)
    add_textbox(slide, 1.1, y + 0.1, 2.2, 0.9, title, font_size=18, color=color, bold=True)
    # divider
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.5), Inches(y + 0.15), Pt(1), Inches(0.8))
    shape.fill.solid(); shape.fill.fore_color.rgb = GRAY_M; shape.line.fill.background()
    add_textbox(slide, 3.8, y + 0.1, 8.2, 0.9, desc, font_size=14, color=GRAY_L)

add_page_number(slide, 11)

# ── Slide 12: Roadmap ─────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_header(slide, "Roadmap", "From research to revenue")

phases = [
    ("Phase A", "Unified Framework", "COMPLETE", "provider_adapter, experiment_runner, pyproject.toml", GREEN),
    ("Phase B", "Scientific Rigor", "COMPLETE", "n=100 × 9 models, P3 neural, 4 families, temp robustness", GREEN),
    ("Phase C", "Paper + Open Source", "NEXT →", "arXiv submission, pip install decay-eval, HuggingFace Space", ACCENT2),
    ("Phase D", "Leaderboard Website", "Q3 2026", "React frontend, FastAPI backend, weekly auto-eval 10+ models", ACCENT),
    ("Phase E", "Paid Services", "Q4 2026", "Pro private eval ($1,499), Enterprise subscriptions ($4,999/mo)", PURPLE),
]

for i, (phase, name, status, detail, color) in enumerate(phases):
    y = 2.2 + i * 1.0
    add_rounded_rect(slide, 0.8, y, 11.5, 0.85, BG_CARD)
    add_circle(slide, 1.0, y + 0.15, 0.5, color)
    add_textbox(slide, 1.0, y + 0.2, 0.5, 0.4, str(i+1), font_size=16, color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 1.8, y + 0.05, 2.0, 0.35, phase, font_size=14, color=color, bold=True)
    add_textbox(slide, 1.8, y + 0.35, 2.0, 0.35, name, font_size=12, color=GRAY_L)
    add_textbox(slide, 4.2, y + 0.2, 2.0, 0.35, status, font_size=16, color=WHITE, bold=True)
    add_textbox(slide, 6.5, y + 0.2, 5.5, 0.35, detail, font_size=12, color=GRAY_L)

add_page_number(slide, 12)

# ── Slide 13: Why Now ─────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_header(slide, "Why Now", "Three converging trends")

trends = [
    ("AI-generated content is flooding the training pipeline",
     "57% of web content projected AI-generated by 2027. Models will inevitably train on their own outputs."),
    ("No standard metric exists for recursive stability",
     "MMLU, HumanEval, Chatbot Arena — all measure single-pass quality. The evaluation matrix has a blank cell."),
    ("LLM companies need pre-release stability diagnostics",
     "Model degradation discovered post-deployment is catastrophic. Pre-release testing creates an insurance market."),
]

for i, (title, desc) in enumerate(trends):
    y = 2.3 + i * 1.6
    add_rounded_rect(slide, 0.8, y, 11.5, 1.35, BG_CARD)
    add_textbox(slide, 1.2, y + 0.2, 10.5, 0.45, f"0{i+1}", font_size=28, color=ACCENT, bold=True)
    add_textbox(slide, 2.0, y + 0.2, 9.8, 0.45, title, font_size=18, color=WHITE, bold=True)
    add_textbox(slide, 2.0, y + 0.7, 9.8, 0.45, desc, font_size=14, color=GRAY_L)

add_page_number(slide, 13)

# ── Slide 14: Closing ─────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_textbox(slide, 1.5, 1.8, 10.3, 1.2, "DecayMonitor", font_size=72, color=WHITE, bold=True)
add_textbox(slide, 1.5, 3.0, 10.3, 0.8, "The missing dimension in LLM evaluation.", font_size=28, color=GRAY_L)
add_thin_line(slide, 1.5, 4.0, 3.5, ACCENT)

add_multiline(slide, 1.5, 4.5, 10.3, 2.0, [
    "Deng Xinhang  ·  Beijing",
    "Independent Research  ·  Constraint AI",
    "",
    "decaymonitor.ai",
], font_size=18, color=GRAY_L, line_spacing=1.5)

add_page_number(slide, 14)

# ── Save ──────────────────────────────────────────────────────────────
output_path = "/Users/dengxinhang/paper/constraint_residual/DecayMonitor_Pitch_Deck.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
