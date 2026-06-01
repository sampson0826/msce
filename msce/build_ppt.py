#!/usr/bin/env python3
"""Build MSCE Product Promotional PPTX."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.oxml.ns import qn
from pptx.chart.data import CategoryChartData
import copy

# ── Constants ──────────────────────────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Colors
BG_DARK    = RGBColor(0x0D, 0x0D, 0x1A)   # Almost black navy
BG_CARD    = RGBColor(0x16, 0x21, 0x3E)   # Navy card
BG_CARD2   = RGBColor(0x1A, 0x27, 0x45)   # Slightly lighter card
CYAN       = RGBColor(0x00, 0xD4, 0xFF)   # #00D4FF
PURPLE     = RGBColor(0x7B, 0x2F, 0xF7)   # #7B2FF7
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY       = RGBColor(0xB0, 0xB0, 0xC8)
LIGHT_GRAY = RGBColor(0x88, 0x88, 0xA0)
GREEN      = RGBColor(0x00, 0xE6, 0x76)
ORANGE     = RGBColor(0xFF, 0x91, 0x00)
RED        = RGBColor(0xFF, 0x52, 0x52)
YELLOW     = RGBColor(0xFF, 0xD6, 0x00)
MAGENTA    = RGBColor(0xE0, 0x40, 0xFB)
CARD_DARK  = RGBColor(0x12, 0x1A, 0x30)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# Remove default slide layouts margins by using blank
blank_layout = prs.slide_layouts[6]  # blank

# ── Helper Functions ───────────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(blank_layout)

def set_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color=None, border_color=None, border_width=None, radius=None):
    """Add a rounded or normal rectangle."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height
    )
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        if border_width:
            shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text="", font_size=Pt(16), color=WHITE,
                bold=False, alignment=PP_ALIGN.LEFT, font_name="Arial", anchor=MSO_ANCHOR.TOP,
                line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(4)
    # Set line spacing
    pPr = p._pPr
    if pPr is None:
        pPr = p._p.get_or_add_pPr()
    lnSpc = pPr.makeelement(qn('a:lnSpc'), {})
    spcPct = lnSpc.makeelement(qn('a:spcPct'), {'val': str(int(line_spacing * 100000))})
    lnSpc.append(spcPct)
    pPr.append(lnSpc)
    # anchor
    txBox.text_frame.paragraphs[0].space_before = Pt(0)
    return txBox

def add_multiline_textbox(slide, left, top, width, height, lines, font_name="Arial", anchor=MSO_ANCHOR.TOP):
    """lines: list of (text, font_size, color, bold, alignment)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, line in enumerate(lines):
        text, size, color, bold, align = line
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = size
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = align
        p.space_after = Pt(2)
    return txBox

def add_slide_number(slide, num):
    add_textbox(slide, Inches(12.4), Inches(7.05), Inches(0.8), Inches(0.35),
                str(num), Pt(11), LIGHT_GRAY, False, PP_ALIGN.RIGHT, "Arial")

def add_accent_line(slide, left, top, width, height=Pt(4), color=CYAN):
    """Add a small accent line (gradient bar)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_gradient_line(slide, left, top, width, height=Pt(4)):
    """Draw a gradient-like effect by stacking two rectangles."""
    half = int(width / 2)
    add_accent_line(slide, left, top, half, height, CYAN)
    add_accent_line(slide, left + half, top, half, height, PURPLE)

def add_card(slide, left, top, width, height, title="", body_lines=None, title_color=CYAN,
             body_color=WHITE, card_color=BG_CARD):
    """Add a card with title and body lines."""
    card = add_rect(slide, left, top, width, height, fill_color=card_color, border_color=RGBColor(0x2A, 0x35, 0x55), border_width=Pt(1))
    y_offset = top + Inches(0.15)
    if title:
        add_textbox(slide, left + Inches(0.2), y_offset, width - Inches(0.4), Inches(0.35),
                    title, Pt(15), title_color, True, PP_ALIGN.LEFT)
        y_offset += Inches(0.35)
    if body_lines:
        for line_text, line_size, line_color, line_bold in body_lines:
            add_textbox(slide, left + Inches(0.2), y_offset, width - Inches(0.4), Inches(0.3),
                        line_text, line_size, line_color, line_bold, PP_ALIGN.LEFT)
            y_offset += Inches(0.28)
    return card

def add_big_number(slide, left, top, number, label, num_color=CYAN, label_color=WHITE):
    """Add a big stat number with label."""
    add_textbox(slide, left, top, Inches(2.5), Inches(1.0),
                number, Pt(48), num_color, True, PP_ALIGN.CENTER)
    add_textbox(slide, left, top + Inches(0.85), Inches(2.5), Inches(0.4),
                label, Pt(14), label_color, False, PP_ALIGN.CENTER)

def add_circle(slide, left, top, size, fill_color=None, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(2)
    return shape

def add_arrow_right(slide, left, top, width, height, color=CYAN):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_chevron(slide, left, top, width, height, color=CYAN):
    shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
set_bg(s, BG_DARK)

# Decorative background - large faded circle
add_circle(s, Inches(8.5), Inches(-1.5), Inches(7), fill_color=RGBColor(0x10, 0x18, 0x30))
add_circle(s, Inches(-2), Inches(4), Inches(5), fill_color=RGBColor(0x0F, 0x14, 0x28))

# Top accent line
add_gradient_line(s, Inches(0), Inches(0), SLIDE_W, Pt(6))

# MSCE logo text
add_textbox(s, Inches(1.5), Inches(1.2), Inches(10), Inches(1.2),
            "MSCE", Pt(72), WHITE, True, PP_ALIGN.LEFT, "Arial")

# Accent bar under MSCE
add_gradient_line(s, Inches(1.5), Inches(2.35), Inches(3.5), Pt(5))

# Chinese name
add_textbox(s, Inches(1.5), Inches(2.55), Inches(10), Inches(0.7),
            "Multi-model Self-Consistency Engine", Pt(24), CYAN, False, PP_ALIGN.LEFT, "Arial")
add_textbox(s, Inches(1.5), Inches(3.1), Inches(10), Inches(0.7),
            "认知对抗引擎", Pt(28), PURPLE, True, PP_ALIGN.LEFT, "Arial")

# Tagline
add_textbox(s, Inches(1.5), Inches(4.0), Inches(10), Inches(0.8),
            '"让 AI 知道自己什么时候不懂"', Pt(22), GRAY, False, PP_ALIGN.LEFT, "Arial")
add_textbox(s, Inches(1.5), Inches(4.55), Inches(10), Inches(0.5),
            '"Let AI know when it doesn\'t know"', Pt(16), LIGHT_GRAY, False, PP_ALIGN.LEFT, "Arial")

# Subtitle
add_textbox(s, Inches(1.5), Inches(5.5), Inches(10), Inches(0.5),
            "v5  |  Open Source (MIT)  |  95% on Chinese Math  |  86.7% on MMLU", Pt(14), GRAY, False, PP_ALIGN.LEFT, "Arial")

add_slide_number(s, 1)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2: THE PROBLEM
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
set_bg(s, BG_DARK)
add_gradient_line(s, Inches(0), Inches(0), SLIDE_W, Pt(5))

add_textbox(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
            "The Problem", Pt(36), WHITE, True, PP_ALIGN.LEFT)
add_textbox(s, Inches(0.8), Inches(1.0), Inches(11), Inches(0.4),
            "AI 自信地给出错误答案 — 而我们无从知晓", Pt(16), GRAY, False, PP_ALIGN.LEFT)

# Three problem cards
card_w = Inches(3.7)
card_h = Inches(4.5)
gap = Inches(0.3)
start_x = Inches(0.8)

# Card 1
c = add_rect(s, start_x, Inches(1.8), card_w, card_h, fill_color=BG_CARD, border_color=RGBColor(0x2A, 0x35, 0x55), border_width=Pt(1))
add_textbox(s, start_x + Inches(0.25), Inches(2.0), card_w - Inches(0.5), Inches(0.4),
            "01", Pt(14), CYAN, True, PP_ALIGN.LEFT)
add_textbox(s, start_x + Inches(0.25), Inches(2.35), card_w - Inches(0.5), Inches(0.5),
            "Hallucination Is Systematic", Pt(18), WHITE, True, PP_ALIGN.LEFT)
add_multiline_textbox(s, start_x + Inches(0.25), Inches(2.9), card_w - Inches(0.5), Inches(3.0), [
    ("Even state-of-the-art models", Pt(13), GRAY, False, PP_ALIGN.LEFT),
    ("confidently produce incorrect", Pt(13), GRAY, False, PP_ALIGN.LEFT),
    ("answers — with no warning signal.", Pt(13), GRAY, False, PP_ALIGN.LEFT),
    ("", Pt(8), GRAY, False, PP_ALIGN.LEFT),
    ("GPT-4o: 70% accuracy on", Pt(13), ORANGE, True, PP_ALIGN.LEFT),
    ("Chinese math benchmark.", Pt(13), ORANGE, True, PP_ALIGN.LEFT),
    ("That's 30% confidently wrong.", Pt(13), ORANGE, True, PP_ALIGN.LEFT),
])

# Card 2
x2 = start_x + card_w + gap
c = add_rect(s, x2, Inches(1.8), card_w, card_h, fill_color=BG_CARD, border_color=RGBColor(0x2A, 0x35, 0x55), border_width=Pt(1))
add_textbox(s, x2 + Inches(0.25), Inches(2.0), card_w - Inches(0.5), Inches(0.4),
            "02", Pt(14), CYAN, True, PP_ALIGN.LEFT)
add_textbox(s, x2 + Inches(0.25), Inches(2.35), card_w - Inches(0.5), Inches(0.5),
            "The 'Overconfident Expert' Trap", Pt(18), WHITE, True, PP_ALIGN.LEFT)
add_multiline_textbox(s, x2 + Inches(0.25), Inches(2.9), card_w - Inches(0.5), Inches(3.0), [
    ("LLMs generate fluent, persuasive", Pt(13), GRAY, False, PP_ALIGN.LEFT),
    ("text regardless of correctness.", Pt(13), GRAY, False, PP_ALIGN.LEFT),
    ("", Pt(8), GRAY, False, PP_ALIGN.LEFT),
    ("No built-in mechanism to", Pt(13), GRAY, False, PP_ALIGN.LEFT),
    ("signal uncertainty or say", Pt(13), GRAY, False, PP_ALIGN.LEFT),
    ('"I don\'t know."', Pt(13), GRAY, False, PP_ALIGN.LEFT),
])

# Card 3
x3 = x2 + card_w + gap
c = add_rect(s, x3, Inches(1.8), card_w, card_h, fill_color=BG_CARD, border_color=RGBColor(0x2A, 0x35, 0x55), border_width=Pt(1))
add_textbox(s, x3 + Inches(0.25), Inches(2.0), card_w - Inches(0.5), Inches(0.4),
            "03", Pt(14), CYAN, True, PP_ALIGN.LEFT)
add_textbox(s, x3 + Inches(0.25), Inches(2.35), card_w - Inches(0.5), Inches(0.5),
            "High-Stakes Consequences", Pt(18), WHITE, True, PP_ALIGN.LEFT)
add_multiline_textbox(s, x3 + Inches(0.25), Inches(2.9), card_w - Inches(0.5), Inches(3.0), [
    ("Legal: wrong answer = liability", Pt(13), GRAY, False, PP_ALIGN.LEFT),
    ("Finance: wrong answer = $ loss", Pt(13), GRAY, False, PP_ALIGN.LEFT),
    ("Medical: wrong answer = harm", Pt(13), GRAY, False, PP_ALIGN.LEFT),
    ("", Pt(8), GRAY, False, PP_ALIGN.LEFT),
    ("You need to know WHEN", Pt(13), ORANGE, True, PP_ALIGN.LEFT),
    ("you cannot trust the AI.", Pt(13), ORANGE, True, PP_ALIGN.LEFT),
])

add_textbox(s, Inches(0.8), Inches(6.6), Inches(11), Inches(0.4),
            "The core problem is not accuracy — it is UNKNOWN inaccuracy.", Pt(14), RED, True, PP_ALIGN.LEFT)
add_slide_number(s, 2)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3: THE INSIGHT
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
set_bg(s, BG_DARK)
add_gradient_line(s, Inches(0), Inches(0), SLIDE_W, Pt(5))

add_textbox(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
            "The Insight", Pt(36), WHITE, True, PP_ALIGN.LEFT)
add_textbox(s, Inches(0.8), Inches(1.0), Inches(11), Inches(0.4),
            "单一模型永远无法「知道自己不知道什么」— MSCE 可以", Pt(16), GRAY, False, PP_ALIGN.LEFT)

# Left: single model limitation
add_rect(s, Inches(0.8), Inches(1.8), Inches(5.8), Inches(5.0), fill_color=BG_CARD, border_color=RGBColor(0x2A, 0x35, 0x55), border_width=Pt(1))
add_textbox(s, Inches(1.1), Inches(2.0), Inches(5.3), Inches(0.5),
            "Single Model Limitation", Pt(22), RED, True, PP_ALIGN.CENTER)
add_circle(s, Inches(2.6), Inches(3.0), Inches(2.0), fill_color=RGBColor(0x20, 0x0A, 0x0A), border_color=RED)
add_textbox(s, Inches(2.6), Inches(3.45), Inches(2.0), Inches(0.6),
            "GPT-4o", Pt(20), RED, True, PP_ALIGN.CENTER)
add_textbox(s, Inches(2.6), Inches(3.85), Inches(2.0), Inches(0.5),
            "70%", Pt(28), RED, True, PP_ALIGN.CENTER)
add_multiline_textbox(s, Inches(1.1), Inches(5.0), Inches(5.3), Inches(1.6), [
    ("A single model has only ONE cognitive", Pt(13), GRAY, False, PP_ALIGN.CENTER),
    ("path through any problem. It cannot", Pt(13), GRAY, False, PP_ALIGN.CENTER),
    ("self-detect when that path is wrong.", Pt(13), GRAY, False, PP_ALIGN.CENTER),
    ("", Pt(6), GRAY, False, PP_ALIGN.CENTER),
    ("This is a fundamental mathematical", Pt(13), ORANGE, True, PP_ALIGN.CENTER),
    ("limitation — not a training problem.", Pt(13), ORANGE, True, PP_ALIGN.CENTER),
])

# Arrow
add_arrow_right(s, Inches(6.7), Inches(3.8), Inches(0.8), Inches(0.5), CYAN)

# Right: MSCE solution
add_rect(s, Inches(7.6), Inches(1.8), Inches(5.2), Inches(5.0), fill_color=RGBColor(0x0A, 0x15, 0x20), border_color=CYAN, border_width=Pt(2))
add_textbox(s, Inches(7.9), Inches(2.0), Inches(4.6), Inches(0.5),
            "MSCE Multi-Angle Approach", Pt(22), GREEN, True, PP_ALIGN.CENTER)

# 6 small circles representing models
model_colors = [CYAN, PURPLE, MAGENTA, ORANGE, GREEN, YELLOW]
for i, mc in enumerate(model_colors):
    row = i // 3
    col = i % 3
    cx = Inches(8.5) + col * Inches(1.4)
    cy = Inches(2.8) + row * Inches(1.3)
    add_circle(s, cx, cy, Inches(0.7), fill_color=RGBColor(0x15, 0x25, 0x35), border_color=mc)
    add_textbox(s, cx, cy + Inches(0.15), Inches(0.7), Inches(0.4),
                f"M{i+1}", Pt(11), mc, True, PP_ALIGN.CENTER)

add_multiline_textbox(s, Inches(7.9), Inches(5.2), Inches(4.6), Inches(1.4), [
    ("6 heterogeneous models attack the", Pt(13), GRAY, False, PP_ALIGN.CENTER),
    ("same problem from different cognitive", Pt(13), GRAY, False, PP_ALIGN.CENTER),
    ("angles. Disagreement reveals uncertainty.", Pt(13), GRAY, False, PP_ALIGN.CENTER),
    ("", Pt(6), GRAY, False, PP_ALIGN.CENTER),
    ("When models disagree,", Pt(14), GREEN, True, PP_ALIGN.CENTER),
    ("MSCE KNOWS it's uncertain.", Pt(14), GREEN, True, PP_ALIGN.CENTER),
])

add_slide_number(s, 3)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4: WHAT IS MSCE
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
set_bg(s, BG_DARK)
add_gradient_line(s, Inches(0), Inches(0), SLIDE_W, Pt(5))

add_textbox(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
            "What is MSCE", Pt(36), WHITE, True, PP_ALIGN.LEFT)
add_textbox(s, Inches(0.8), Inches(1.0), Inches(11), Inches(0.4),
            "6 Models. 6 Cognitive Strategies. 1 AI Courtroom.", Pt(18), CYAN, True, PP_ALIGN.LEFT)

# Big number: 6
add_textbox(s, Inches(0.8), Inches(1.8), Inches(3), Inches(1.5),
            "6", Pt(120), CYAN, True, PP_ALIGN.LEFT)
add_textbox(s, Inches(0.8), Inches(3.2), Inches(3), Inches(0.4),
            "Heterogeneous Models", Pt(16), WHITE, True, PP_ALIGN.LEFT)

# Details
details = [
    ("GPT-4o", "OpenAI's flagship, broad reasoning"),
    ("DeepSeek-V3", "MoE architecture, math strength"),
    ("Claude 4", "Long-context, careful reasoning"),
    ("Qwen3-235B", "Multilingual, strong coding"),
    ("Gemini 2.5 Pro", "Google's multimodal power"),
    ("Grok-3", "xAI's real-time knowledge"),
]
for i, (name, desc) in enumerate(details):
    y = Inches(3.8) + i * Inches(0.55)
    col = CYAN if i < 3 else PURPLE
    add_circle(s, Inches(1.0), y + Inches(0.02), Inches(0.35), fill_color=col)
    add_textbox(s, Inches(1.5), y, Inches(2.5), Inches(0.25), name, Pt(13), WHITE, True)
    add_textbox(s, Inches(1.5), y + Inches(0.22), Inches(2.5), Inches(0.2), desc, Pt(10), GRAY, False)

# Right side: the courtroom concept
add_rect(s, Inches(5.0), Inches(1.8), Inches(7.8), Inches(5.0), fill_color=BG_CARD, border_color=RGBColor(0x2A, 0x35, 0x55), border_width=Pt(1))

add_textbox(s, Inches(5.3), Inches(2.0), Inches(7.2), Inches(0.5),
            "The Multi-Angle Adversarial Process", Pt(20), WHITE, True, PP_ALIGN.CENTER)

# Flow: Question → 6 models → Adversarial elimination → Verdict
steps = [
    ("1 Question", CYAN),
    ("6 Answers", PURPLE),
    ("Adversarial\nElimination", ORANGE),
    ("Judge\nVerdict", GREEN),
]
sw = Inches(1.6)
for i, (label, col) in enumerate(steps):
    sx = Inches(5.5) + i * Inches(1.85)
    add_rect(s, sx, Inches(2.7), sw, Inches(1.2), fill_color=RGBColor(0x15, 0x25, 0x35), border_color=col, border_width=Pt(2))
    add_textbox(s, sx, Inches(2.9), sw, Inches(0.8), label, Pt(13), col, True, PP_ALIGN.CENTER)
    if i < 3:
        add_arrow_right(s, sx + sw + Inches(0.05), Inches(3.05), Inches(0.2), Inches(0.3), GRAY)

# Output
add_textbox(s, Inches(5.3), Inches(4.3), Inches(7.2), Inches(0.4),
            "Output: Answer + Confidence Score + Disagreement Metrics", Pt(14), GREEN, True, PP_ALIGN.CENTER)

# Three key metrics
metrics = [
    ("Confidence", "How sure is the\nfinal answer?", CYAN),
    ("Disagreement", "How much did\nmodels differ?", ORANGE),
    ("I Don't Know", "Signal when AI\nshould not be trusted", RED),
]
for i, (title, desc, col) in enumerate(metrics):
    mx = Inches(5.5) + i * Inches(2.5)
    add_textbox(s, mx, Inches(4.9), Inches(2.2), Inches(0.3), title, Pt(15), col, True, PP_ALIGN.CENTER)
    add_textbox(s, mx, Inches(5.2), Inches(2.2), Inches(0.6), desc, Pt(11), GRAY, False, PP_ALIGN.CENTER)

# Tagline
add_textbox(s, Inches(5.3), Inches(6.0), Inches(7.2), Inches(0.5),
            "Not better AI — SMARTER certainty about your AI.", Pt(14), CYAN, True, PP_ALIGN.CENTER)
add_slide_number(s, 4)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5: HOW IT WORKS — Architecture Flow
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
set_bg(s, BG_DARK)
add_gradient_line(s, Inches(0), Inches(0), SLIDE_W, Pt(5))

add_textbox(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
            "How It Works", Pt(36), WHITE, True, PP_ALIGN.LEFT)
add_textbox(s, Inches(0.8), Inches(1.0), Inches(11), Inches(0.4),
            "Architecture: From Question to Trusted Answer", Pt(16), GRAY, False, PP_ALIGN.LEFT)

# Stage boxes
stages = [
    ("Input\nQuestion", "User submits\na query", CYAN),
    ("6× Generator\nModels", "Parallel reasoning\nby 6 different AIs", PURPLE),
    ("Answer\nCollection", "Gather all\n6 responses", MAGENTA),
    ("Adversarial\nElimination", "Cross-validate,\neliminate outliers", ORANGE),
    ("Consensus\nJudge", "LLM judge evaluates\nremaining answers", YELLOW),
    ("Output: Answer\n+ Confidence", "Trusted answer\nwith metrics", GREEN),
]

for i, (title, desc, col) in enumerate(stages):
    sx = Inches(0.4) + i * Inches(2.1)
    sy = Inches(1.9)
    bw = Inches(1.9)
    bh = Inches(2.3)
    add_rect(s, sx, sy, bw, bh, fill_color=BG_CARD, border_color=col, border_width=Pt(2))
    add_textbox(s, sx + Inches(0.1), sy + Inches(0.1), bw - Inches(0.2), Inches(0.7),
                title, Pt(14), col, True, PP_ALIGN.CENTER)
    add_textbox(s, sx + Inches(0.1), sy + Inches(0.9), bw - Inches(0.2), Inches(1.0),
                desc, Pt(11), GRAY, False, PP_ALIGN.CENTER)
    if i < 5:
        ax = sx + bw + Inches(0.02)
        add_arrow_right(s, ax, sy + Inches(1.0), Inches(0.16), Inches(0.25), GRAY)

# Bottom detail cards
detail_cards = [
    ("Parallel Execution", "All 6 models run simultaneously,\nminimizing latency overhead", CYAN),
    ("Semantic Clustering", "Answers grouped by semantic\nsimilarity, not surface form", PURPLE),
    ("LLM-as-Judge", "A separate judge model evaluates\nconsensus to produce final verdict", ORANGE),
    ("Confidence Calibration", "Agreement level mapped to\ncalibrated confidence score", GREEN),
]
for i, (title, desc, col) in enumerate(detail_cards):
    dx = Inches(0.4) + i * Inches(3.25)
    add_rect(s, dx, Inches(4.8), Inches(3.0), Inches(2.2), fill_color=RGBColor(0x10, 0x18, 0x2A), border_color=RGBColor(0x25, 0x30, 0x4A), border_width=Pt(1))
    add_textbox(s, dx + Inches(0.15), Inches(4.95), Inches(2.7), Inches(0.3),
                title, Pt(14), col, True, PP_ALIGN.LEFT)
    add_textbox(s, dx + Inches(0.15), Inches(5.3), Inches(2.7), Inches(1.3),
                desc, Pt(11), GRAY, False, PP_ALIGN.LEFT)

add_slide_number(s, 5)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6: THE AI COURTROOM
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
set_bg(s, BG_DARK)
add_gradient_line(s, Inches(0), Inches(0), SLIDE_W, Pt(5))

add_textbox(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
            'The "AI Courtroom"', Pt(36), WHITE, True, PP_ALIGN.LEFT)
add_textbox(s, Inches(0.8), Inches(1.0), Inches(11), Inches(0.4),
            "6 AI models debate. One Judge. One Verdict.", Pt(16), CYAN, True, PP_ALIGN.LEFT)

# Courtroom layout: 6 models as "jury" on sides, judge at top center
# Judge box (top center)
add_rect(s, Inches(4.5), Inches(1.7), Inches(4.3), Inches(1.8), fill_color=RGBColor(0x1A, 0x10, 0x05), border_color=YELLOW, border_width=Pt(3))
add_textbox(s, Inches(4.5), Inches(1.8), Inches(4.3), Inches(0.5),
            "JUDGE (LLM)", Pt(22), YELLOW, True, PP_ALIGN.CENTER)
add_textbox(s, Inches(4.5), Inches(2.3), Inches(4.3), Inches(0.4),
            "Evaluates consensus, cross-references evidence,", Pt(12), GRAY, False, PP_ALIGN.CENTER)
add_textbox(s, Inches(4.5), Inches(2.6), Inches(4.3), Inches(0.4),
            "produces final answer with confidence score", Pt(12), GRAY, False, PP_ALIGN.CENTER)

# 6 model boxes: 3 on left, 3 on right
model_info = [
    ("M1: GPT-4o", "Broad reasoning", CYAN),
    ("M2: DeepSeek-V3", "Math & logic", PURPLE),
    ("M3: Claude 4", "Careful analysis", MAGENTA),
    ("M4: Qwen3-235B", "Multilingual", ORANGE),
    ("M5: Gemini 2.5 Pro", "Multimodal", GREEN),
    ("M6: Grok-3", "Real-time data", YELLOW),
]

for i, (name, desc, col) in enumerate(model_info):
    if i < 3:
        mx = Inches(0.8)
        my = Inches(1.7) + i * Inches(1.6)
    else:
        mx = Inches(10.0)
        my = Inches(1.7) + (i - 3) * Inches(1.6)
    add_rect(s, mx, my, Inches(2.8), Inches(1.2), fill_color=BG_CARD, border_color=col, border_width=Pt(2))
    add_textbox(s, mx + Inches(0.1), my + Inches(0.15), Inches(2.6), Inches(0.35),
                name, Pt(14), col, True, PP_ALIGN.CENTER)
    add_textbox(s, mx + Inches(0.1), my + Inches(0.5), Inches(2.6), Inches(0.3),
                desc, Pt(11), GRAY, False, PP_ALIGN.CENTER)

    # Arrow pointing to center (toward judge)
    if i < 3:
        ax = mx + Inches(2.8)
        add_arrow_right(s, ax, my + Inches(0.4), Inches(0.8), Inches(0.3), GRAY)
    else:
        ax = mx - Inches(0.8)
        # left arrow
        shape = s.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, ax, my + Inches(0.4), Inches(0.8), Inches(0.3))
        shape.fill.solid()
        shape.fill.fore_color.rgb = GRAY
        shape.line.fill.background()

# Bottom: deliberation flow
add_rect(s, Inches(1.5), Inches(5.5), Inches(10.3), Inches(1.5), fill_color=RGBColor(0x0A, 0x12, 0x22), border_color=RGBColor(0x25, 0x30, 0x4A), border_width=Pt(1))

delib_steps = [
    ("1. Parallel Reasoning", "Each model independently\nsolves the problem"),
    ("2. Semantic Clustering", "Answers grouped by\nsemantic similarity"),
    ("3. Adversarial Challenge", "Models cross-examine\neach other's outputs"),
    ("4. Consensus Extraction", "Judge identifies the\nbest-supported answer"),
]
for i, (title, desc) in enumerate(delib_steps):
    dx = Inches(1.8) + i * Inches(2.7)
    add_textbox(s, dx, Inches(5.65), Inches(2.4), Inches(0.3),
                title, Pt(13), CYAN, True, PP_ALIGN.LEFT)
    add_textbox(s, dx, Inches(5.95), Inches(2.4), Inches(0.6),
                desc, Pt(10), GRAY, False, PP_ALIGN.LEFT)
    if i < 3:
        add_arrow_right(s, dx + Inches(2.35), Inches(6.0), Inches(0.2), Inches(0.22), GRAY)

add_slide_number(s, 6)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7: BENCHMARK RESULTS — Chinese Math
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
set_bg(s, BG_DARK)
add_gradient_line(s, Inches(0), Inches(0), SLIDE_W, Pt(5))

add_textbox(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
            "Benchmark: Chinese Math (20 Questions)", Pt(36), WHITE, True, PP_ALIGN.LEFT)
add_textbox(s, Inches(0.8), Inches(1.0), Inches(11), Inches(0.4),
            "v5 Results — MSCE dominates all single-model baselines", Pt(16), GRAY, False, PP_ALIGN.LEFT)

# Bar chart
chart_data = CategoryChartData()
chart_data.categories = ['GPT-4o', 'Claude 4', 'Gemini\n2.5 Pro', 'DeepSeek\nV3', 'Qwen3\n235B', 'MSCE v5']
chart_data.add_series('Accuracy (%)', (70, 75, 80, 90, 85, 95))

chart_frame = s.shapes.add_chart(
    1,  # XL_CHART_TYPE.COLUMN_CLUSTERED
    Inches(0.8), Inches(1.6), Inches(7.5), Inches(5.0),
    chart_data
)
chart = chart_frame.chart
chart.has_legend = False
chart.has_title = False

# Style the chart
plot = chart.plots[0]
series = plot.series[0]
# Set gradient-like colors for bars
colors = [RGBColor(0x88, 0x88, 0xA0), RGBColor(0x88, 0x88, 0xA0), RGBColor(0x88, 0x88, 0xA0),
          RGBColor(0x88, 0x88, 0xA0), RGBColor(0x88, 0x88, 0xA0), RGBColor(0x00, 0xD4, 0xFF)]
for idx, color in enumerate(colors):
    point = series.points[idx]
    point.format.fill.solid()
    point.format.fill.fore_color.rgb = color

# Chart area styling
chart_frame.chart.font.color.rgb = WHITE
chart_frame.chart.font.size = Pt(11)
chart_frame.chart.font.name = 'Arial'

# Value axis
value_axis = chart.value_axis
value_axis.minimum_scale = 0
value_axis.maximum_scale = 100
value_axis.major_gridlines.format.line.color.rgb = RGBColor(0x30, 0x30, 0x50)
value_axis.format.line.color.rgb = GRAY
value_axis.tick_labels.font.color.rgb = GRAY
value_axis.tick_labels.font.size = Pt(10)

# Category axis
cat_axis = chart.category_axis
cat_axis.tick_labels.font.color.rgb = GRAY
cat_axis.tick_labels.font.size = Pt(10)
cat_axis.format.line.color.rgb = RGBColor(0x30, 0x30, 0x50)

# Chart border removal not needed with dark theme

# Right side: big stats
add_rect(s, Inches(8.8), Inches(1.6), Inches(4.0), Inches(5.3), fill_color=BG_CARD, border_color=RGBColor(0x2A, 0x35, 0x55), border_width=Pt(1))

add_big_number(s, Inches(8.8), Inches(1.9), "95%", "MSCE v5 Accuracy", GREEN)
add_textbox(s, Inches(8.8), Inches(3.4), Inches(4.0), Inches(0.3),
            "+25% vs GPT-4o", Pt(14), CYAN, True, PP_ALIGN.CENTER)
add_textbox(s, Inches(8.8), Inches(3.7), Inches(4.0), Inches(0.3),
            "+5% vs DeepSeek-V3", Pt(14), PURPLE, True, PP_ALIGN.CENTER)

add_textbox(s, Inches(8.8), Inches(4.3), Inches(4.0), Inches(0.3),
            "---", Pt(14), GRAY, False, PP_ALIGN.CENTER)

# Additional stats
stats_lines = [
    ("19/20 correct", "MSCE v5", GREEN),
    ("18/20 correct", "DeepSeek-V3", PURPLE),
    ("14/20 correct", "GPT-4o", RED),
]
for i, (score, model, col) in enumerate(stats_lines):
    y = Inches(4.7) + i * Inches(0.4)
    dot = add_circle(s, Inches(9.1), y + Inches(0.05), Inches(0.2), fill_color=col)
    add_textbox(s, Inches(9.4), y, Inches(1.5), Inches(0.25), score, Pt(13), WHITE, True)
    add_textbox(s, Inches(10.8), y, Inches(1.8), Inches(0.25), model, Pt(11), GRAY, False)

add_textbox(s, Inches(8.8), Inches(6.1), Inches(4.0), Inches(0.4),
            "20-question Chinese math benchmark\n(addition, multiplication, algebra, logic)", Pt(10), LIGHT_GRAY, False, PP_ALIGN.CENTER)

add_slide_number(s, 7)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8: MMLU PILOT
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
set_bg(s, BG_DARK)
add_gradient_line(s, Inches(0), Inches(0), SLIDE_W, Pt(5))

add_textbox(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
            "MMLU Pilot (30 Questions)", Pt(36), WHITE, True, PP_ALIGN.LEFT)
add_textbox(s, Inches(0.8), Inches(1.0), Inches(11), Inches(0.4),
            "MSCE 86.7% — outperforming every single model", Pt(16), GRAY, False, PP_ALIGN.LEFT)

# Bar chart
chart_data2 = CategoryChartData()
chart_data2.categories = ['GPT-4o', 'Claude 4', 'Gemini\n2.5 Pro', 'DeepSeek\nV3', 'Qwen3\n235B', 'MSCE v5']
chart_data2.add_series('Accuracy (%)', (56.7, 73.3, 76.7, 83.3, 80.0, 86.7))

chart_frame2 = s.shapes.add_chart(
    1,  # COLUMN_CLUSTERED
    Inches(0.8), Inches(1.6), Inches(7.5), Inches(4.8),
    chart_data2
)
chart2 = chart_frame2.chart
chart2.has_legend = False
chart2.has_title = False

plot2 = chart2.plots[0]
series2 = plot2.series[0]
colors2 = [RGBColor(0x88, 0x88, 0xA0)] * 5 + [RGBColor(0x7B, 0x2F, 0xF7)]
for idx, color in enumerate(colors2):
    point = series2.points[idx]
    point.format.fill.solid()
    point.format.fill.fore_color.rgb = color

chart2.font.color.rgb = WHITE
chart2.font.size = Pt(11)
chart2.font.name = 'Arial'

va2 = chart2.value_axis
va2.minimum_scale = 0
va2.maximum_scale = 100
va2.major_gridlines.format.line.color.rgb = RGBColor(0x30, 0x30, 0x50)
va2.format.line.color.rgb = GRAY
va2.tick_labels.font.color.rgb = GRAY
va2.tick_labels.font.size = Pt(10)

ca2 = chart2.category_axis
ca2.tick_labels.font.color.rgb = GRAY
ca2.tick_labels.font.size = Pt(10)
ca2.format.line.color.rgb = RGBColor(0x30, 0x30, 0x50)

# Right side stats
add_rect(s, Inches(8.8), Inches(1.6), Inches(4.0), Inches(4.8), fill_color=BG_CARD, border_color=RGBColor(0x2A, 0x35, 0x55), border_width=Pt(1))

add_big_number(s, Inches(8.8), Inches(1.9), "86.7%", "MSCE v5 MMLU", PURPLE)
add_textbox(s, Inches(8.8), Inches(3.4), Inches(4.0), Inches(0.3),
            "+30% vs GPT-4o", Pt(14), CYAN, True, PP_ALIGN.CENTER)
add_textbox(s, Inches(8.8), Inches(3.7), Inches(4.0), Inches(0.3),
            "+3.4% vs DeepSeek-V3", Pt(14), PURPLE, True, PP_ALIGN.CENTER)

# MMLU subjects
add_textbox(s, Inches(8.8), Inches(4.3), Inches(4.0), Inches(0.3),
            "MMLU Subjects Tested:", Pt(13), WHITE, True, PP_ALIGN.LEFT)
subjects = ["High School Math", "College Physics", "Professional Law",
            "Econometrics", "Moral Scenarios", "Computer Science"]
for i, subj in enumerate(subjects):
    add_textbox(s, Inches(9.0), Inches(4.6) + i * Inches(0.25), Inches(3.6), Inches(0.22),
                f"  *  {subj}", Pt(10), GRAY, False)

add_textbox(s, Inches(8.8), Inches(6.1), Inches(4.0), Inches(0.3),
            "30-question pilot across 6 MMLU subjects", Pt(10), LIGHT_GRAY, False, PP_ALIGN.CENTER)

add_slide_number(s, 8)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9: UNCERTAINTY QUANTIFICATION
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
set_bg(s, BG_DARK)
add_gradient_line(s, Inches(0), Inches(0), SLIDE_W, Pt(5))

add_textbox(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
            "Uncertainty Quantification: The Secret Sauce", Pt(36), WHITE, True, PP_ALIGN.LEFT)
add_textbox(s, Inches(0.8), Inches(1.0), Inches(11), Inches(0.4),
            "Not just an answer — a calibrated signal of when not to trust", Pt(16), CYAN, True, PP_ALIGN.LEFT)

# Three metric cards
metrics_data = [
    ("Confidence Score", "0.0 — 1.0", "Agreement ratio among winning\nanswer cluster vs. all responses.\nCalibrated to match actual accuracy.", CYAN,
     ["* High confidence ( >0.8): 94% correct", "* Medium (0.5-0.8): review recommended", "* Low (<0.5): do NOT trust"]),
    ("Disagreement Index", "0.0 — 1.0", "Entropy across model responses.\nMeasures how much models diverge\non this specific question.", ORANGE,
     ["* Low disagreement: models agree", "* High disagreement: human review needed", "* Triggers 'I don't know' flag"]),
    ("I Don't Know Signal", "Binary + Score", "When confidence < threshold OR\ndisagreement > threshold, MSCE\nrefuses to give an answer.", RED,
     ["* Active uncertainty detection", "* Prevents silent failures", "* Routes to human expert"]),
]

for i, (title, range_str, desc, col, bullets) in enumerate(metrics_data):
    mx = Inches(0.8) + i * Inches(4.1)
    mw = Inches(3.8)
    add_rect(s, mx, Inches(1.7), mw, Inches(5.3), fill_color=BG_CARD, border_color=col, border_width=Pt(2))

    # Title
    add_textbox(s, mx + Inches(0.2), Inches(1.85), mw - Inches(0.4), Inches(0.35),
                title, Pt(20), col, True, PP_ALIGN.LEFT)
    # Range badge
    add_rect(s, mx + Inches(0.2), Inches(2.2), Inches(1.5), Inches(0.32), fill_color=col)
    add_textbox(s, mx + Inches(0.2), Inches(2.2), Inches(1.5), Inches(0.32),
                range_str, Pt(11), BG_DARK, True, PP_ALIGN.CENTER)
    # Description
    add_textbox(s, mx + Inches(0.2), Inches(2.7), mw - Inches(0.4), Inches(1.2),
                desc, Pt(12), GRAY, False)
    # Bullets
    for j, bullet in enumerate(bullets):
        add_textbox(s, mx + Inches(0.3), Inches(4.0) + j * Inches(0.35), mw - Inches(0.6), Inches(0.3),
                    bullet, Pt(11), WHITE, False)

add_slide_number(s, 9)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10: WHY THIS MATTERS — High-Stakes AI
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
set_bg(s, BG_DARK)
add_gradient_line(s, Inches(0), Inches(0), SLIDE_W, Pt(5))

add_textbox(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
            "Why This Matters", Pt(36), WHITE, True, PP_ALIGN.LEFT)
add_textbox(s, Inches(0.8), Inches(1.0), Inches(11), Inches(0.5),
            "Single-model AI fails silently. MSCE catches what they miss.", Pt(16), RED, True, PP_ALIGN.LEFT)

# Three industry cards
industries = [
    ("LEGAL", "Wrong Answer = Liability", RED,
     ["Contract analysis errors lead to", "malpractice claims", "",
      "MSCE confidence score prevents", "deployment of uncertain outputs"]),
    ("FINANCE", "Wrong Answer = Dollar Loss", ORANGE,
     ["Trading signals, risk assessment,", "compliance — errors cost millions", "",
      "MSCE disagreement detection flags", "high-risk recommendations"]),
    ("MEDICAL", "Wrong Answer = Harm", RED,
     ["Diagnostic suggestions, drug", "interactions — lives at stake", "",
      "MSCE 'I don't know' signal ensures", "uncertain cases go to human experts"]),
]

for i, (label, subtitle, col, bullets) in enumerate(industries):
    ix = Inches(0.8) + i * Inches(4.1)
    iw = Inches(3.8)
    add_rect(s, ix, Inches(1.8), iw, Inches(4.8), fill_color=BG_CARD, border_color=col, border_width=Pt(2))

    # Icon circle
    add_circle(s, ix + Inches(1.3), Inches(1.95), Inches(1.2), fill_color=RGBColor(0x20, 0x0A, 0x0A), border_color=col)
    add_textbox(s, ix + Inches(1.3), Inches(2.3), Inches(1.2), Inches(0.5),
                label, Pt(14), col, True, PP_ALIGN.CENTER)

    add_textbox(s, ix + Inches(0.2), Inches(3.3), iw - Inches(0.4), Inches(0.35),
                subtitle, Pt(16), col, True, PP_ALIGN.CENTER)
    for j, bullet in enumerate(bullets):
        add_textbox(s, ix + Inches(0.3), Inches(3.75) + j * Inches(0.28), iw - Inches(0.6), Inches(0.25),
                    bullet, Pt(12), GRAY if bullet else WHITE, False)

# Bottom banner
add_rect(s, Inches(0.8), Inches(6.85), Inches(11.7), Inches(0.45), fill_color=RGBColor(0x1A, 0x08, 0x08), border_color=RED, border_width=Pt(1))
add_textbox(s, Inches(0.8), Inches(6.85), Inches(11.7), Inches(0.45),
            "In high-stakes domains, knowing when NOT to trust AI is more valuable than higher average accuracy.", Pt(13), RED, True, PP_ALIGN.CENTER)

add_slide_number(s, 10)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11: FOR DEVELOPERS
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
set_bg(s, BG_DARK)
add_gradient_line(s, Inches(0), Inches(0), SLIDE_W, Pt(5))

add_textbox(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
            "For Developers", Pt(36), WHITE, True, PP_ALIGN.LEFT)
add_textbox(s, Inches(0.8), Inches(1.0), Inches(11), Inches(0.4),
            "Simple API. Multiple SDKs. Docker One-Click.", Pt(16), CYAN, True, PP_ALIGN.LEFT)

# Code snippet card
add_rect(s, Inches(0.8), Inches(1.7), Inches(7.5), Inches(5.2), fill_color=RGBColor(0x0D, 0x14, 0x24), border_color=RGBColor(0x25, 0x30, 0x4A), border_width=Pt(1))

# Fake code block
code_lines = [
    ("from msce import Engine", CYAN),
    ("", WHITE),
    ("# Initialize with your API keys", GRAY),
    ("engine = Engine(", WHITE),
    ("    models=['gpt-4o', 'deepseek-v3',", WHITE),
    ("            'claude-4', 'qwen3-235b',", WHITE),
    ("            'gemini-2.5-pro', 'grok-3'],", WHITE),
    ("    judge='claude-4'", WHITE),
    (")", WHITE),
    ("", WHITE),
    ("# Ask a question", GRAY),
    ("result = engine.ask(", WHITE),
    ('    "What is the integral of x^2 * ln(x)?"', WHITE),
    (")", WHITE),
    ("", WHITE),
    ("print(result.answer)        # 'x^3/3 * ...'", WHITE),
    ("print(result.confidence)    # 0.94", GREEN),
    ("print(result.disagreement)  # 0.12", ORANGE),
    ("print(result.models_agreed) # 5/6", PURPLE),
]

for i, (text, col) in enumerate(code_lines):
    add_textbox(s, Inches(1.1), Inches(1.85) + i * Inches(0.28), Inches(7.0), Inches(0.26),
                text, Pt(11), col, False, PP_ALIGN.LEFT, "Courier New")

# Right side: features
add_rect(s, Inches(8.8), Inches(1.7), Inches(4.0), Inches(5.2), fill_color=BG_CARD, border_color=RGBColor(0x2A, 0x35, 0x55), border_width=Pt(1))

sdk_features = [
    ("Python SDK", "pip install msce", CYAN),
    ("JavaScript/TS SDK", "npm install msce-sdk", PURPLE),
    ("REST API", "api.msce.dev/v1/ask", MAGENTA),
    ("LangChain Integration", "from langchain_msce import ...", ORANGE),
    ("Docker Deploy", "docker compose up -d", GREEN),
    ("Streaming Support", "SSE real-time results", CYAN),
]

for i, (title, cmd, col) in enumerate(sdk_features):
    fy = Inches(1.9) + i * Inches(0.78)
    add_rect(s, Inches(9.0), fy, Inches(3.6), Inches(0.65), fill_color=RGBColor(0x12, 0x1A, 0x30), border_color=RGBColor(0x25, 0x30, 0x4A), border_width=Pt(1))
    add_textbox(s, Inches(9.1), fy + Inches(0.02), Inches(3.4), Inches(0.28),
                title, Pt(13), col, True, PP_ALIGN.LEFT)
    add_textbox(s, Inches(9.1), fy + Inches(0.32), Inches(3.4), Inches(0.25),
                cmd, Pt(11), GRAY, False, PP_ALIGN.LEFT, "Courier New")

add_slide_number(s, 11)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12: OPEN SOURCE
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
set_bg(s, BG_DARK)
add_gradient_line(s, Inches(0), Inches(0), SLIDE_W, Pt(5))

add_textbox(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
            "Open Source & Transparent", Pt(36), WHITE, True, PP_ALIGN.LEFT)
add_textbox(s, Inches(0.8), Inches(1.0), Inches(11), Inches(0.4),
            "MIT License. Your API Keys. Full Control.", Pt(16), GREEN, True, PP_ALIGN.LEFT)

# Three principle cards
principles = [
    ("MIT License", "100% open source.", "Use it, modify it, build on it.\nNo restrictions. No vendor lock-in.\nNo hidden fees.", CYAN),
    ("You Own Your Keys", "Bring your own API keys.", "You pay your model providers\ndirectly. MSCE never sees\nyour keys or your data.", PURPLE),
    ("Transparent Architecture", "Every line of code visible.", "No black boxes. No magic.\nUnderstand exactly how\nconfidence scores are computed.", MAGENTA),
]

for i, (title, subtitle, desc, col) in enumerate(principles):
    px = Inches(0.8) + i * Inches(4.1)
    pw = Inches(3.8)
    add_rect(s, px, Inches(1.8), pw, Inches(3.8), fill_color=BG_CARD, border_color=col, border_width=Pt(2))

    add_circle(s, px + Inches(1.3), Inches(1.95), Inches(1.2), fill_color=RGBColor(0x10, 0x15, 0x25), border_color=col)
    # Checkmark icon (simplified as text)
    add_textbox(s, px + Inches(1.3), Inches(2.25), Inches(1.2), Inches(0.5),
                "OPEN", Pt(14), col, True, PP_ALIGN.CENTER)
    add_textbox(s, px + Inches(0.15), Inches(3.3), pw - Inches(0.3), Inches(0.35),
                title, Pt(18), col, True, PP_ALIGN.CENTER)
    add_textbox(s, px + Inches(0.15), Inches(3.7), pw - Inches(0.3), Inches(0.3),
                subtitle, Pt(13), WHITE, True, PP_ALIGN.CENTER)
    add_textbox(s, px + Inches(0.15), Inches(4.1), pw - Inches(0.3), Inches(1.2),
                desc, Pt(11), GRAY, False)

# Cline model reference
add_rect(s, Inches(0.8), Inches(6.0), Inches(11.7), Inches(1.2), fill_color=RGBColor(0x0A, 0x15, 0x20), border_color=CYAN, border_width=Pt(1))
add_textbox(s, Inches(1.0), Inches(6.1), Inches(3.0), Inches(0.3),
            "Business Model: The Cline Approach", Pt(16), CYAN, True)
add_textbox(s, Inches(1.0), Inches(6.4), Inches(7.0), Inches(0.5),
            "Open source tool + user-pays-their-own-API-keys. Proven model — Cline raised $32M.", Pt(13), GRAY, False)
add_textbox(s, Inches(11.2), Inches(6.4), Inches(1.3), Inches(0.5),
            "$32M", Pt(22), GREEN, True, PP_ALIGN.RIGHT)

add_slide_number(s, 12)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 13: COMPARISON TABLE
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
set_bg(s, BG_DARK)
add_gradient_line(s, Inches(0), Inches(0), SLIDE_W, Pt(5))

add_textbox(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
            "Comparison: MSCE vs The World", Pt(36), WHITE, True, PP_ALIGN.LEFT)
add_textbox(s, Inches(0.8), Inches(1.0), Inches(11), Inches(0.4),
            "No other system quantifies its own uncertainty", Pt(16), CYAN, True, PP_ALIGN.LEFT)

# Build comparison table with shapes
table_left = Inches(0.8)
table_top = Inches(1.7)
col_w = [Inches(2.0), Inches(2.5), Inches(2.5), Inches(2.7), Inches(2.7)]
row_h = Inches(0.5)
headers = ["Feature", "GPT-4o", "DeepSeek-V3", "Claude 4", "MSCE v5"]
rows_data = [
    ["Math Accuracy (CN)", "70%", "90%", "75%", "95%"],
    ["MMLU Accuracy", "56.7%", "83.3%", "73.3%", "86.7%"],
    ["Confidence Score", "No", "No", "No", "Yes"],
    ["Uncertainty Detection", "No", "No", "No", "Yes"],
    ["Multi-Model Consensus", "No", "No", "No", "Yes"],
    ["Disagreement Metrics", "No", "No", "No", "Yes"],
    ["'I Don't Know' Signal", "No", "No", "No", "Yes"],
    ["Open Source", "No", "Partial", "No", "MIT Full"],
    ["User-Owned API Keys", "N/A", "N/A", "N/A", "Yes"],
]

# Header row
hx = table_left
for j, (cw, hdr) in enumerate(zip(col_w, headers)):
    is_last = (j == len(headers) - 1)
    hdr_color = GREEN if is_last else CYAN
    add_rect(s, hx, table_top, cw, row_h, fill_color=RGBColor(0x18, 0x25, 0x40), border_color=RGBColor(0x2A, 0x35, 0x55), border_width=Pt(1))
    add_textbox(s, hx, table_top + Inches(0.05), cw, Inches(0.4),
                hdr, Pt(13), hdr_color, True, PP_ALIGN.CENTER)
    hx += cw

# Data rows
for i, row in enumerate(rows_data):
    ry = table_top + (i + 1) * row_h
    rbg = BG_CARD if i % 2 == 0 else RGBColor(0x14, 0x1E, 0x38)
    hx = table_left
    for j, (cw, cell) in enumerate(zip(col_w, row)):
        is_last = (j == len(row) - 1)
        is_bad = cell in ("No", "N/A")
        cell_color = GREEN if is_last else (RED if is_bad else (GREEN if cell == "Yes" else WHITE))
        add_rect(s, hx, ry, cw, row_h, fill_color=rbg, border_color=RGBColor(0x20, 0x2A, 0x4A), border_width=Pt(0.5))
        add_textbox(s, hx, ry + Inches(0.08), cw, Inches(0.35),
                    cell, Pt(12), cell_color, is_last, PP_ALIGN.CENTER)
        hx += cw

add_slide_number(s, 13)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 14: ROADMAP
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
set_bg(s, BG_DARK)
add_gradient_line(s, Inches(0), Inches(0), SLIDE_W, Pt(5))

add_textbox(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
            "Roadmap", Pt(36), WHITE, True, PP_ALIGN.LEFT)
add_textbox(s, Inches(0.8), Inches(1.0), Inches(11), Inches(0.4),
            "Building toward universal AI certainty", Pt(16), GRAY, False, PP_ALIGN.LEFT)

# Timeline
phases = [
    ("NOW", "Q2 2026", "Math & Logic", CYAN, [
        "Chinese math benchmark: 95%",
        "MMLU pilot: 86.7%",
        "Python SDK released",
        "Docker one-click deploy",
        "GitHub public (MIT)",
    ]),
    ("NEXT", "Q3 2026", "Domain Expansion", PURPLE, [
        "Legal reasoning optimization",
        "Financial analysis models",
        "Medical QA specialization",
        "LangChain + LlamaIndex",
        "Confidence calibration v2",
    ]),
    ("FUTURE", "Q4 2026+", "Enterprise", MAGENTA, [
        "Enterprise hosted version",
        "Custom model integration",
        "On-premise deployment",
        "SOC 2 compliance",
        "SLAs for production use",
    ]),
]

for i, (badge, timeline, title, col, items) in enumerate(phases):
    px = Inches(0.8) + i * Inches(4.1)
    pw = Inches(3.8)

    # Badge
    add_rect(s, px, Inches(1.7), pw, Inches(0.5), fill_color=col)
    add_textbox(s, px, Inches(1.72), pw, Inches(0.45),
                f"{badge} — {timeline}", Pt(15), BG_DARK, True, PP_ALIGN.CENTER)

    # Card
    add_rect(s, px, Inches(2.2), pw, Inches(4.6), fill_color=BG_CARD, border_color=col, border_width=Pt(1))
    add_textbox(s, px + Inches(0.15), Inches(2.4), pw - Inches(0.3), Inches(0.4),
                title, Pt(20), col, True, PP_ALIGN.CENTER)
    add_gradient_line(s, px + Inches(1.0), Inches(2.85), pw - Inches(2.0), Pt(3))

    for j, item in enumerate(items):
        is_done = item.endswith("%") and badge == "NOW"
        ic = GREEN if is_done else WHITE
        prefix = "*" if is_done else "~"
        add_textbox(s, px + Inches(0.3), Inches(3.1) + j * Inches(0.4), pw - Inches(0.6), Inches(0.35),
                    f"  {prefix}  {item}", Pt(12), ic, False)

    # Connecting arrow
    if i < 2:
        add_arrow_right(s, px + pw + Inches(0.05), Inches(1.8), Inches(0.2), Inches(0.3), GRAY)

add_slide_number(s, 14)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 15: GET STARTED
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
set_bg(s, BG_DARK)
add_gradient_line(s, Inches(0), Inches(0), SLIDE_W, Pt(5))

add_textbox(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
            "Get Started in 60 Seconds", Pt(36), WHITE, True, PP_ALIGN.LEFT)
add_textbox(s, Inches(0.8), Inches(1.0), Inches(11), Inches(0.4),
            "Open source. Self-hosted. Production ready.", Pt(16), GREEN, True, PP_ALIGN.LEFT)

# Four quick-start cards
qs_cards = [
    ("GitHub", "github.com/msce-ai/msce", "Star, fork, contribute.\nReport issues, suggest features.", CYAN, "git clone"),
    ("Docker", "docker compose up -d", "One command. Everything\nrunning in seconds.", PURPLE, "docker"),
    ("API Endpoint", "POST /v1/ask", "REST API with JSON.\nStreaming via SSE.", MAGENTA, "curl"),
    ("Community", "Discord + GitHub Discussions", "Get help, share use cases,\nbuild together.", ORANGE, "chat"),
]

for i, (title, cmd, desc, col, icon) in enumerate(qs_cards):
    qx = Inches(0.8) + i * Inches(3.2)
    qw = Inches(2.95)
    add_rect(s, qx, Inches(1.7), qw, Inches(3.5), fill_color=BG_CARD, border_color=col, border_width=Pt(2))

    # Icon
    add_circle(s, qx + Inches(0.95), Inches(1.85), Inches(1.0), fill_color=RGBColor(0x10, 0x15, 0x25), border_color=col)
    add_textbox(s, qx + Inches(0.95), Inches(2.1), Inches(1.0), Inches(0.5),
                icon, Pt(11), col, True, PP_ALIGN.CENTER)

    add_textbox(s, qx + Inches(0.1), Inches(3.0), qw - Inches(0.2), Inches(0.3),
                title, Pt(16), col, True, PP_ALIGN.CENTER)
    add_textbox(s, qx + Inches(0.1), Inches(3.3), qw - Inches(0.2), Inches(0.5),
                cmd, Pt(11), WHITE, False, PP_ALIGN.CENTER, "Courier New")
    add_textbox(s, qx + Inches(0.1), Inches(3.9), qw - Inches(0.2), Inches(1.0),
                desc, Pt(11), GRAY, False, PP_ALIGN.CENTER)

# Bottom install command banner
add_rect(s, Inches(0.8), Inches(5.6), Inches(11.7), Inches(1.3), fill_color=RGBColor(0x0D, 0x14, 0x24), border_color=RGBColor(0x25, 0x30, 0x4A), border_width=Pt(1))
add_textbox(s, Inches(1.0), Inches(5.7), Inches(2.0), Inches(0.3),
            "Quick Install:", Pt(14), CYAN, True)
add_textbox(s, Inches(1.0), Inches(6.05), Inches(11.0), Inches(0.4),
            "$ pip install msce", Pt(22), GREEN, False, PP_ALIGN.LEFT, "Courier New")
add_textbox(s, Inches(1.0), Inches(6.5), Inches(11.0), Inches(0.3),
            "$ msce ask 'What is the capital of France?'", Pt(16), GRAY, False, PP_ALIGN.LEFT, "Courier New")

add_slide_number(s, 15)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 16: CLOSING
# ═══════════════════════════════════════════════════════════════════════════
s = add_slide()
set_bg(s, BG_DARK)

# Decorative background elements
add_circle(s, Inches(9), Inches(-2), Inches(8), fill_color=RGBColor(0x10, 0x18, 0x30))
add_circle(s, Inches(-3), Inches(3), Inches(6), fill_color=RGBColor(0x0F, 0x14, 0x28))

add_gradient_line(s, Inches(0), Inches(0), SLIDE_W, Pt(6))

add_textbox(s, Inches(1.0), Inches(1.5), Inches(11.3), Inches(1.0),
            "MSCE", Pt(72), WHITE, True, PP_ALIGN.CENTER)
add_gradient_line(s, Inches(4.5), Inches(2.5), Inches(4.3), Pt(5))

add_textbox(s, Inches(1.0), Inches(2.8), Inches(11.3), Inches(0.7),
            "Multi-model Self-Consistency Engine", Pt(22), CYAN, False, PP_ALIGN.CENTER)
add_textbox(s, Inches(1.0), Inches(3.5), Inches(11.3), Inches(0.7),
            "认知对抗引擎", Pt(26), PURPLE, True, PP_ALIGN.CENTER)
add_textbox(s, Inches(1.0), Inches(4.3), Inches(11.3), Inches(0.6),
            "Know When to Trust AI.", Pt(30), WHITE, True, PP_ALIGN.CENTER)
add_textbox(s, Inches(1.0), Inches(4.9), Inches(11.3), Inches(0.5),
            "让AI知道自己什么时候不懂。", Pt(20), GRAY, False, PP_ALIGN.CENTER)

# Stats row
add_rect(s, Inches(1.5), Inches(5.6), Inches(10.3), Inches(0.01), fill_color=RGBColor(0x30, 0x30, 0x55))
stats_final = [
    ("95%", "Chinese Math", CYAN),
    ("86.7%", "MMLU", PURPLE),
    ("6", "Models", MAGENTA),
    ("MIT", "License", GREEN),
]
for i, (val, label, col) in enumerate(stats_final):
    sx = Inches(2.0) + i * Inches(2.7)
    add_textbox(s, sx, Inches(5.75), Inches(2.0), Inches(0.5),
                val, Pt(36), col, True, PP_ALIGN.CENTER)
    add_textbox(s, sx, Inches(6.25), Inches(2.0), Inches(0.3),
                label, Pt(13), GRAY, False, PP_ALIGN.CENTER)

add_textbox(s, Inches(1.0), Inches(6.85), Inches(11.3), Inches(0.4),
            "github.com/msce-ai/msce  |  MIT License  |  Open Source", Pt(12), LIGHT_GRAY, False, PP_ALIGN.CENTER)

add_slide_number(s, 16)

# ── Save ───────────────────────────────────────────────────────────────────
output_path = "/Users/dengxinhang/paper/constraint_residual/msce/MSCE_产品宣传.pptx"
prs.save(output_path)
print(f"PPTX saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
